from flask import Flask, request, render_template, send_file, redirect, url_for
import pytesseract
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches
import re
import os

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
REPORT_PATH = "report.pptx"

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Define function to interpret image and generate caption
def interpret_image(image_path):
    # Preprocess the image
    image = Image.open(image_path)
    image = ImageOps.grayscale(image)  # Convert to grayscale
    extracted_text = pytesseract.image_to_string(image)

    print("Extracted Text:", extracted_text)  # Debugging step

    # Adjust regex to match numbers followed by "Max" or "Min"
    match_max = re.search(r"(\d+\.\d+)\s*Max", extracted_text, re.IGNORECASE)
    match_min = re.search(r"(\d+\.\d+)\s*Min", extracted_text, re.IGNORECASE)

    max_value = match_max.group(1) if match_max else None
    min_value = match_min.group(1) if match_min else None

    # Classify image and generate caption
    if max_value and min_value:
        caption = f"The max stress is {max_value}.\nThe min stress is {min_value}."
        image_type = "Stress"
    elif max_value:
        caption = f"The max deformation is {max_value}."
        image_type = "Deformation"
    else:
        caption = "Unable to classify the image."
        image_type = "Unknown"

    return image_type, caption

# Define function to create PowerPoint report
def create_ppt(image_path, caption, image_type):
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 6 corresponds to a blank layout
    slide = prs.slides.add_slide(slide_layout)

    heading = "Stress Analysis" if image_type == "Stress" else "Deformation Analysis"
    # title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(9), Inches(1))
    # title_frame = title_box.text_frame
    # title_frame.text = heading
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = heading
    title_frame.paragraphs[0].font.size = Inches(0.7)

    image_left = Inches(0.5)  # Position the image on the left
    image_top = Inches(2)
    image_width = Inches(4.5)
    slide.shapes.add_picture(image_path, image_left, image_top, image_width, Inches(4))

    caption_left = Inches(5.5)  # Position the caption on the right
    caption_top = Inches(3)
    caption_box = slide.shapes.add_textbox(caption_left, caption_top, Inches(4), Inches(2))
    caption_frame = caption_box.text_frame
    caption_frame.text = caption

    prs.save(REPORT_PATH)

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        file = request.files["image"]
        if file:
            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            file.save(file_path)

            # Process the image and generate the report
            image_type, caption = interpret_image(file_path)
            create_ppt(file_path, caption, image_type)

            return redirect(url_for("download"))

    return render_template("index.html")

@app.route("/download")
def download():
    return send_file(REPORT_PATH, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
